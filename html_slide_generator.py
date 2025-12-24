"""
PDF/Image template-based slide generator.
Uses Canva PDF or image template and overlays text/images using PIL.
Supports both PDF and image templates (JPG, PNG, etc.)
"""
import os
from typing import Dict, Optional
import io
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import hashlib
import img2pdf
from collections import Counter
import numpy as np
from functools import lru_cache

# PPTX support for editable Canva designs
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Fix for reportlab compatibility
try:
    hashlib.md5(b'test', usedforsecurity=False)
except TypeError:
    _original_md5 = hashlib.md5
    def _patched_md5(data=None, usedforsecurity=True):
        return _original_md5(data) if data is not None else _original_md5()
    hashlib.md5 = _patched_md5

class HTMLSlideGenerator:
    def __init__(self):
        from config import Config
        
        # Get the directory where this script is located (works on Render too)
        script_dir = os.path.dirname(os.path.abspath(__file__))
        # On Render: script_dir = /opt/render/project/src, project_root = /opt/render/project
        # Locally: script_dir = /path/to/slauson-automation, project_root = /path/to/slauson-automation
        # Check if we're in a 'src' subdirectory (Render) or directly in project root
        if os.path.basename(script_dir) == 'src':
            project_root = os.path.dirname(script_dir)
        else:
            project_root = script_dir
        
        print(f"DEBUG: script_dir={script_dir}, project_root={project_root}")
        
        # Slide template path - check config first, then try default locations
        self.template_path = getattr(Config, 'SLIDE_TEMPLATE_PATH', None) if hasattr(Config, 'SLIDE_TEMPLATE_PATH') else None
        if not self.template_path or not os.path.exists(self.template_path):
            # Try default locations (project root first, then templates folder, then script dir)
            default_template_paths = [
                # Project root (where user added the file)
                os.path.join(project_root, 'SLAUSON&CO.Template.pdf'),
                os.path.join(project_root, 'SLAUSON&CO.template'),
                os.path.join(project_root, 'SLAUSON&CO.Template'),
                # Templates folder
                os.path.join(project_root, 'templates', 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'templates', 'SLAUSON&CO.Template.pdf'),
                'templates/SLAUSON&CO.Template.pdf',
                # Script directory
                os.path.join(script_dir, 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'SLAUSON&CO.template'),
                # Current working directory
                'SLAUSON&CO.Template.pdf',
                'SLAUSON&CO.template',
                # Generic template names
                os.path.join(project_root, 'templates', 'template.pdf'),
                os.path.join(script_dir, 'templates', 'template.pdf'),
                'templates/template.pdf',
                'templates/template.png',
                'templates/template.jpg',
            ]
            print(f"DEBUG: Checking {len(default_template_paths)} template paths...")
            for path in default_template_paths:
                if path and os.path.exists(path):
                    self.template_path = path
                    print(f"✓ Found template at: {self.template_path}")
                    break
                else:
                    print(f"  - Not found: {path}")
        
        # Map template path - check config first, then try default locations
        self.map_template_path = getattr(Config, 'MAP_TEMPLATE_PATH', None) if hasattr(Config, 'MAP_TEMPLATE_PATH') else None
        if not self.map_template_path or not os.path.exists(self.map_template_path):
            # Try default locations (relative to script, then project root, then current working directory)
            default_map_paths = [
                os.path.join(script_dir, 'templates', 'map_template.pdf'),
                os.path.join(project_root, 'templates', 'map_template.pdf'),
                'templates/map_template.pdf',
                os.path.join(script_dir, 'templates', 'SLAUSON&CO. (1).pdf'),
                os.path.join(project_root, 'templates', 'SLAUSON&CO. (1).pdf'),
                'templates/SLAUSON&CO. (1).pdf',
            ]
            for path in default_map_paths:
                if path and os.path.exists(path):
                    self.map_template_path = path
                    print(f"✓ Found map template at: {self.map_template_path}")
                    break

    def _load_font(self, size: int, bold: bool = False):
        """
        Load a font with robust fallbacks that work on Render.
        Tries common system fonts, then DejaVu (available in most containers), then default.
        """
        font_candidates = []
        # Preferred bundled fonts on Linux containers (Render uses these)
        if bold:
            font_candidates.append("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf")
            font_candidates.append("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf")
        font_candidates.append("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
        font_candidates.append("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf")
        # Try to find fonts via fontconfig (common on Linux)
        try:
            import subprocess
            result = subprocess.run(['fc-list'], capture_output=True, text=True, timeout=1)
            if result.returncode == 0:
                # Look for DejaVu or Liberation in font list
                for line in result.stdout.split('\n'):
                    if 'DejaVu' in line or 'Liberation' in line:
                        font_path = line.split(':')[0] if ':' in line else None
                        if font_path and os.path.exists(font_path):
                            if bold and 'Bold' in font_path:
                                font_candidates.insert(0, font_path)
                            elif not bold and 'Bold' not in font_path:
                                font_candidates.insert(0, font_path)
        except Exception:
            pass  # Fontconfig not available, continue with hardcoded paths
        
        # macOS fonts (for local development only)
        if bold:
            font_candidates.append("/System/Library/Fonts/Helvetica-Bold.ttf")
            font_candidates.append("/System/Library/Fonts/Arial Bold.ttf")
        font_candidates.append("/System/Library/Fonts/Helvetica.ttc")
        font_candidates.append("/System/Library/Fonts/Arial.ttf")

        for path in font_candidates:
            try:
                if os.path.exists(path) or not os.path.isabs(path):
                    # Only check existence for absolute paths
                    return ImageFont.truetype(path, size)
            except (OSError, IOError, Exception) as e:
                # Silently continue to next candidate
                continue
        
        # Final fallback - always works
        print(f"Warning: Could not load custom font, using default font for size {size}")
        return ImageFont.load_default()
    
    def _pdf_to_image(self, pdf_path: str) -> Image.Image:
        """Convert first page of PDF to PIL Image."""
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=300, first_page=1, last_page=1)
            if images:
                return images[0].convert('RGBA')
        except ImportError:
            # Fallback: try PyPDF2 + extract images
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(pdf_path)
                page = reader.pages[0]
                # Try to extract images from PDF
                if '/XObject' in page['/Resources']:
                    xObject = page['/Resources']['/XObject'].get_object()
                    for obj in xObject:
                        if xObject[obj]['/Subtype'] == '/Image':
                            # Extract image data
                            data = xObject[obj].get_data()
                            img = Image.open(io.BytesIO(data))
                            return img.convert('RGBA')
                # If no images found, render PDF as image using alternative method
                raise ImportError("pdf2image required for PDF templates")
            except Exception as e:
                raise ImportError(
                    f"Could not convert PDF to image. Please install pdf2image: pip install pdf2image. "
                    f"Error: {e}"
                )
        except Exception as e:
            raise Exception(f"Failed to convert PDF to image: {e}")
    
    def _get_dominant_color(self, img: Image.Image, region: tuple, exclude_colors: list = None) -> tuple:
        """Get dominant color in a region."""
        x, y, w, h = region
        region_img = img.crop((x, y, x + w, y + h))
        pixels = list(region_img.getdata())
        
        if exclude_colors:
            pixels = [p for p in pixels if p[:3] not in exclude_colors]
        
        if not pixels:
            return (42, 42, 42)  # Default dark grey
        
        color_counts = Counter([p[:3] for p in pixels])
        return color_counts.most_common(1)[0][0]
    
    def _get_text_color_from_template(self, template: Image.Image, x: int, y: int, width: int, height: int) -> tuple:
        """Extract text color from template."""
        region = template.crop((x, y, x + width, y + height))
        pixels = list(region.getdata())
        
        bg_color = Counter([p[:3] for p in pixels]).most_common(1)[0][0]
        bg_brightness = sum(bg_color) / 3
        
        text_pixels = []
        for p in pixels:
            rgb = p[:3] if len(p) >= 3 else p
            brightness = sum(rgb) / 3
            if abs(brightness - bg_brightness) > 40:
                text_pixels.append(rgb)
        
        if text_pixels:
            return Counter(text_pixels).most_common(1)[0][0]
        
        return (255, 140, 0) if y < 200 else (255, 255, 255)
    
    def _remove_background_manual(self, img: Image.Image, tol: int = 38, feather: int = 2) -> Image.Image:
        """
        Flood-fill from the border to remove background (RGBA alpha=0),
        using multi-corner background samples + auto tolerance fallback.
        """
        try:
            import numpy as np
        except ImportError:
            print("   Warning: numpy not available for manual background removal")
            return img.convert("RGBA")

        img = img.convert("RGBA")
        arr = np.array(img)
        rgb = arr[..., :3].astype(np.int16)
        alpha = arr[..., 3].astype(np.uint8)
        H, W = rgb.shape[:2]

        corner_size = max(12, min(H, W) // 18)

        # Compute per-corner medians (better for gradients / non-uniform bg)
        tl = np.median(rgb[0:corner_size, 0:corner_size].reshape(-1, 3), axis=0)
        tr = np.median(rgb[0:corner_size, W - corner_size:W].reshape(-1, 3), axis=0)
        bl = np.median(rgb[H - corner_size:H, 0:corner_size].reshape(-1, 3), axis=0)
        br = np.median(rgb[H - corner_size:H, W - corner_size:W].reshape(-1, 3), axis=0)
        bgs = np.stack([tl, tr, bl, br], axis=0).astype(np.int16)

        def compute_close_mask(t):
            diffs = rgb[None, ...] - bgs[:, None, None, :]
            dist2 = (diffs * diffs).sum(axis=3)  # (4,H,W)
            dist = np.sqrt(dist2.min(axis=0))    # (H,W)
            return dist <= t

        tol_candidates = [tol, 45, 55, 65, 75] if tol < 45 else [tol, tol + 10, tol + 20]
        best = None

        from collections import deque

        for t in tol_candidates:
            close = compute_close_mask(t)

            bg_mask = np.zeros((H, W), dtype=bool)
            q = deque()

            def push(y, x):
                if 0 <= y < H and 0 <= x < W and close[y, x] and not bg_mask[y, x]:
                    bg_mask[y, x] = True
                    q.append((y, x))

            # seed edges
            for x in range(W):
                push(0, x); push(H - 1, x)
            for y in range(H):
                push(y, 0); push(y, W - 1)

            # 8-neighborhood flood fill
            while q:
                y, x = q.popleft()
                for dy in (-1, 0, 1):
                    for dx in (-1, 0, 1):
                        if dy == 0 and dx == 0:
                            continue
                        push(y + dy, x + dx)

            removed_frac = bg_mask.mean()
            if 0.05 <= removed_frac <= 0.80:
                best = (t, bg_mask, removed_frac)
                break

            if best is None or abs(removed_frac - 0.30) < abs(best[2] - 0.30):
                best = (t, bg_mask, removed_frac)

        t, bg_mask, removed_frac = best
        print(f"   Manual BG removal: tol={t}, removed={removed_frac:.1%}")

        new_alpha = alpha.copy()
        new_alpha[bg_mask] = 0

        out = arr.copy()
        out[..., 3] = new_alpha
        out_img = Image.fromarray(out, "RGBA")

        if feather and feather > 0:
            a = out_img.split()[-1].filter(ImageFilter.GaussianBlur(radius=min(feather, 2)))
            out_img.putalpha(a)

        return out_img

    def _remove_background_gray(self, img: Image.Image, tol: int = 18, feather: int = 2) -> Image.Image:
        """
        Flood-fill background removal for GRAYSCALE headshots.
        Uses luminance-only distance with border sampling and auto-tuned tolerance.
        """
        try:
            import numpy as np
            except ImportError:
            print("   Warning: numpy not available for gray background removal")
            return img.convert("RGBA")

        img = img.convert("RGBA")
        arr = np.array(img)
        H, W = arr.shape[:2]

        # Luminance (even if already grayscale, this is safe)
        lum = (0.299 * arr[..., 0] + 0.587 * arr[..., 1] + 0.114 * arr[..., 2]).astype(np.float32)
        alpha = arr[..., 3].astype(np.uint8)

        # Border-strip background samples in luminance (safer than corners for tight crops)
        bs = max(8, min(H, W) // 25)
        border = np.concatenate([
            lum[:bs, :].ravel(),
            lum[-bs:, :].ravel(),
            lum[:, :bs].ravel(),
            lum[:, -bs:].ravel(),
        ])
        bg = np.median(border)
        dist = np.abs(lum - bg)

        from collections import deque

        def flood(t):
            close = dist <= t
            bg_mask = np.zeros((H, W), dtype=bool)
            q = deque()

            def push(y, x):
                if 0 <= y < H and 0 <= x < W and close[y, x] and not bg_mask[y, x]:
                    bg_mask[y, x] = True
                    q.append((y, x))

            for x in range(W):
                push(0, x); push(H - 1, x)
            for y in range(H):
                push(y, 0); push(y, W - 1)

            while q:
                y, x = q.popleft()
                for dy in (-1, 0, 1):
                    for dx in (-1, 0, 1):
                        if dy == 0 and dx == 0:
                            continue
                        push(y + dy, x + dx)

            return bg_mask

        tol_candidates = [8, 10, 12, 14, 16, 18]
        best_mask = None
        best_score = None
        best_t = tol_candidates[0]

        for t in tol_candidates:
            mask = flood(t)
            removed = mask.mean()
            score = abs(removed - 0.30)  # target ~30% removal
            if best_score is None or score < best_score:
                best_score = score
                best_mask = mask
                best_t = t
            if 0.08 <= removed <= 0.80:
                break

        print(f"   Gray BG removal: tol={best_t}, removed={best_mask.mean():.1%}")

        new_alpha = alpha.copy()
        new_alpha[best_mask] = 0

        out = arr.copy()
        out[..., 3] = new_alpha
        out_img = Image.fromarray(out, "RGBA")

        if feather and feather > 0:
            a = out_img.split()[-1].filter(ImageFilter.GaussianBlur(radius=min(feather, 2)))
            out_img.putalpha(a)

        return out_img

    def _alpha_stats(self, im: Image.Image):
        """Return (opaque_frac, transparent_frac, mean_alpha)."""
        a = np.array(im.split()[-1], dtype=np.uint8)
        opaque = (a > 200).mean()
        transparent = (a < 10).mean()
        return float(opaque), float(transparent), float(a.mean())

    def _remove_bg_rembg(self, rgba_img: Image.Image) -> Optional[Image.Image]:
        """
        Local ML segmentation fallback (rembg / U2Net). Requires: pip install rembg onnxruntime
        """
        try:
            from rembg import remove
            buf = io.BytesIO()
            rgba_img.save(buf, format="PNG")
            out = remove(buf.getvalue())
            out_img = Image.open(io.BytesIO(out)).convert("RGBA")
            out_img.load()
            return out_img
        except Exception as e:
            print(f"   rembg failed: {e}")
            return None

    def _remove_bg_best_effort(self, path: str, use_api: bool) -> Image.Image:
        """
        Best-effort background removal:
        1) remove.bg API if available
        2) local rembg segmentation
        3) conservative grayscale flood-fill as last resort
        """
        img = Image.open(path).convert("RGBA")
        img.load()

        # Downscale before anything expensive
        max_size = 1500
        if max(img.size) > max_size:
            img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)

        # 1) remove.bg API (if available)
        if use_api:
            try:
                from image_processor import ImageProcessor
                print(f"   Removing background via API for {path} ...")
                b = ImageProcessor.remove_background(path)
                if b:
                    api_img = Image.open(io.BytesIO(b)).convert("RGBA")
                    api_img.load()
                    o, t, ma = self._alpha_stats(api_img)
                    print(f"   API alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
                    if o > 0.10 and t > 0.10:
                        return api_img
            except Exception as e:
                print(f"   API removal failed: {e}")

        # 2) local rembg
        rembg_img = self._remove_bg_rembg(img)
        if rembg_img is not None:
            o, t, ma = self._alpha_stats(rembg_img)
            print(f"   rembg alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
            if o > 0.10 and t > 0.10:
                return rembg_img

        # 3) Last resort: conservative flood-fill (no feather)
        ff = self._remove_background_gray(img, tol=10, feather=0)
        o, t, ma = self._alpha_stats(ff)
        print(f"   floodfill alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
        return ff
    
    def _detect_orange_us_bbox(self, img: Image.Image):
        """
        Detect the orange US map outline bbox in the template.
        Restricts search to the TOP-RIGHT region to avoid the orange sidebar.
        Returns (x0, y0, w, h). Falls back if it fails.
        """
        try:
            W, H = img.size
            arr = np.array(img.convert("RGB"))
            
            # --- ROI: where the map is on your template ---
            # tune these if needed, but this matches your layout (map is top-right)
            x_start = int(W * 0.55)
            x_end   = int(W * 0.99)
            y_start = int(H * 0.02)
            y_end   = int(H * 0.45)
            
            roi = arr[y_start:y_end, x_start:x_end]
            r, g, b = roi[..., 0], roi[..., 1], roi[..., 2]
            
            # Orange outline is bright orange lines, not a solid block:
            # loosen thresholds a bit so thin lines still count
            orange_mask = (r > 160) & (g > 80) & (g < 210) & (b < 140)
            
            ys, xs = np.where(orange_mask)
            if len(xs) < 300:  # thin outlines -> fewer pixels; don't require 2000
                print(f"   Warning: Only {len(xs)} orange pixels in ROI, using fallback bbox")
                return (1250, 110, 620, 450)
            
            # bbox in ROI coords
            x0r, x1r = int(xs.min()), int(xs.max())
            y0r, y1r = int(ys.min()), int(ys.max())
            
            # Convert ROI bbox back to full-image coords
            x0 = x_start + x0r
            x1 = x_start + x1r
            y0 = y_start + y0r
            y1 = y_start + y1r
            
            pad = 10
            x0 = max(0, x0 - pad)
            y0 = max(0, y0 - pad)
            x1 = min(W - 1, x1 + pad)
            y1 = min(H - 1, y1 + pad)
            
            print(f"   Detected MAP bbox (ROI): ({x0}, {y0}, {x1-x0}, {y1-y0})")
            return (x0, y0, x1 - x0, y1 - y0)
        except Exception as e:
            print(f"   Warning: Error detecting orange bbox: {e}, using fallback")
            return (1250, 110, 620, 450)
    
    @lru_cache(maxsize=256)
    def _geocode_city(self, city: str):
        """
        Geocode city name to (lat, lon) using geopy.
        Returns (lat, lon) or None.
        Cached to avoid rate limits and improve performance.
        """
        try:
            from geopy.geocoders import Nominatim
            geolocator = Nominatim(user_agent="slide_pin_placer")
            # Bias to US
            loc = geolocator.geocode(f"{city}, USA", country_codes="us", timeout=3)
            if not loc:
                return None
            print(f"   Geocoded '{city}' to ({loc.latitude}, {loc.longitude})")
            return (loc.latitude, loc.longitude)
        except ImportError:
            print(f"   Warning: geopy not installed, cannot geocode '{city}'")
            return None
        except Exception as e:
            print(f"   Warning: Geocoding failed for '{city}': {e}")
            return None
    
    def _fallback_latlon(self, city: str):
        """
        Fallback lat/lon for common cities and all US state capitals when geopy isn't available.
        Uses real coordinates (no nudges) except for LA which is intentionally shifted.
        Note: AK/HI coordinates will clamp to contiguous US bounds unless handled specially.
        """
        # Parse city name - handle both "City, State" and "City State" formats
        parts = city.lower().split(',')
        city_part = parts[0].strip()
        state_part = parts[1].strip() if len(parts) > 1 else ""
        
        LATLON = {
            # Keep these as-is per your request (even though LA is intentionally shifted)
            "los angeles": (34.0522, -118.5),  # intentionally shifted
            "san francisco": (37.7749, -122.4194),
            "miami": (25.7617, -80.1918),
            "new york": (40.7128, -74.0060),
            "new york city": (40.7128, -74.0060),
            
            # Major cities
            "seattle": (47.6062, -122.3321),
            "boston": (42.3601, -71.0589),
            "chicago": (41.8781, -87.6298),
            "new orleans": (29.9511, -90.0715),
            
            # US State Capitals (standard coords)
            "montgomery": (32.3668, -86.3000),        # AL
            "juneau": (58.3019, -134.4197),           # AK  (needs special handling on contiguous map)
            "phoenix": (33.4484, -112.0740),          # AZ
            "little rock": (34.7465, -92.2896),       # AR
            "sacramento": (38.5816, -121.4944),       # CA
            "denver": (39.7392, -104.9903),           # CO
            "hartford": (41.7658, -72.6734),          # CT
            "dover": (39.1582, -75.5244),             # DE
            "tallahassee": (30.4383, -84.2807),       # FL
            "atlanta": (33.7490, -84.3880),           # GA
            "honolulu": (21.3069, -157.8583),         # HI  (needs special handling on contiguous map)
            "boise": (43.6150, -116.2023),            # ID
            "springfield": (39.7817, -89.6501),       # IL (capital)
            "indianapolis": (39.7684, -86.1581),      # IN
            "des moines": (41.5868, -93.6250),        # IA
            "topeka": (39.0473, -95.6752),            # KS
            "frankfort": (38.2009, -84.8733),         # KY
            "baton rouge": (30.4515, -91.1871),       # LA
            "augusta": (44.3106, -69.7795),           # ME (capital)  <-- note ambiguity with Augusta, GA
            "annapolis": (38.9784, -76.4922),         # MD
            "lansing": (42.7325, -84.5555),           # MI
            "saint paul": (44.9537, -93.0900),        # MN
            "st. paul": (44.9537, -93.0900),          # MN alias
            "jackson": (32.2988, -90.1848),           # MS
            "jefferson city": (38.5767, -92.1735),    # MO
            "helena": (46.5891, -112.0391),           # MT
            "lincoln": (40.8136, -96.7026),           # NE
            "carson city": (39.1638, -119.7674),      # NV
            "concord": (43.2081, -71.5376),           # NH
            "trenton": (40.2171, -74.7429),           # NJ
            "santa fe": (35.6870, -105.9378),         # NM
            "albany": (42.6526, -73.7562),            # NY (capital)
            "raleigh": (35.7796, -78.6382),           # NC
            "bismarck": (46.8083, -100.7837),         # ND
            "columbus": (39.9612, -82.9988),          # OH
            "oklahoma city": (35.4676, -97.5164),     # OK
            "salem": (44.9429, -123.0351),            # OR
            "harrisburg": (40.2732, -76.8867),        # PA
            "providence": (41.8240, -71.4128),        # RI
            "columbia": (34.0007, -81.0348),          # SC (capital)
            "pierre": (44.3683, -100.3510),           # SD
            "nashville": (36.1627, -86.7816),         # TN
            "austin": (30.2672, -97.7431),            # TX
            "salt lake city": (40.7608, -111.8910),   # UT
            "montpelier": (44.2601, -72.5754),        # VT
            "richmond": (37.5407, -77.4360),          # VA
            "olympia": (47.0379, -122.9007),          # WA
            "charleston": (38.3498, -81.6326),        # WV (capital) -- ambiguous with Charleston, SC city
            "madison": (43.0748, -89.3848),           # WI
            "cheyenne": (41.1400, -104.8202),         # WY
            
            # Optional disambiguation aliases (HIGHLY recommended)
            "augusta me": (44.3106, -69.7795),
            "augusta ga": (33.4735, -82.0105),
            "charleston wv": (38.3498, -81.6326),
            "charleston sc": (32.7765, -79.9311),
        }
        
        # Try disambiguation key first (e.g., "augusta me", "charleston wv")
        if state_part:
            # Map common state abbreviations to full state names for disambiguation
            state_map = {
                "me": "me", "maine": "me",
                "ga": "ga", "georgia": "ga",
                "wv": "wv", "west virginia": "wv",
                "sc": "sc", "south carolina": "sc"
            }
            state_key = state_map.get(state_part, "")
            if state_key:
                disambiguation_key = f"{city_part} {state_key}"
                if disambiguation_key in LATLON:
                    return LATLON[disambiguation_key]
        
        # Fall back to city name only
        key = city_part
        return LATLON.get(key)
    
    
    def _latlon_to_map_xy(self, lat: float, lon: float, map_x: int, map_y: int, map_w: int, map_h: int):
        """
        Improved lat/lon to map coordinate conversion with calibration support.
        Uses known city positions to calibrate the projection.
        """
        # More accurate contiguous US bounds
        lon_min, lon_max = -124.5, -67.0  # Adjusted for better accuracy
        lat_min, lat_max = 25.0, 49.0
        
        # Adjust padding to match the actual map borders in your template
        # You may need to fine-tune these based on your specific template
        PAD_L = 0.03  # Left padding
        PAD_R = 0.05  # Right padding  
        PAD_T = 0.08  # Top padding (increased to lower West Coast cities)
        PAD_B = 0.08  # Bottom padding
        
        inner_x = map_x + int(map_w * PAD_L)
        inner_y = map_y + int(map_h * PAD_T)
        inner_w = int(map_w * (1.0 - PAD_L - PAD_R))
        inner_h = int(map_h * (1.0 - PAD_T - PAD_B))
        
        # Clamp to bounds
        lon = max(lon_min, min(lon_max, lon))
        lat = max(lat_min, min(lat_max, lat))
        
        # Equirectangular projection
        x_norm = (lon - lon_min) / (lon_max - lon_min)
        y_norm = 1.0 - (lat - lat_min) / (lat_max - lat_min)
        
        x = inner_x + int(x_norm * inner_w)
        y = inner_y + int(y_norm * inner_h)
        
        return x, y
    
    def calibrate_map_projection(self, template_path: str):
        """
        Helper to calibrate map projection by testing known cities.
        Run this once to find the best padding values.
        """
        template = self._pdf_to_image(template_path) if template_path.endswith('.pdf') else Image.open(template_path)
        template.load()  # Load fully before resize to avoid memory issues
        template = template.resize((1920, 1080), Image.Resampling.LANCZOS)
        
        map_x, map_y, map_w, map_h = self._detect_orange_us_bbox(template)
        
        # Test cities with known positions
        test_cities = {
            'Los Angeles': (34.0522, -118.2437),
            'New York': (40.7128, -74.0060),
            'Chicago': (41.8781, -87.6298),
            'Miami': (25.7617, -80.1918),
            'Seattle': (47.6062, -122.3321),
        }
        
        print(f"\nMap bounds: x={map_x}, y={map_y}, w={map_w}, h={map_h}")
        print("\nTesting pin positions:")
        for city, (lat, lon) in test_cities.items():
            x, y = self._latlon_to_map_xy(lat, lon, map_x, map_y, map_w, map_h)
            print(f"{city:15} -> ({x}, {y})")
        
        # Draw test pins on template
        test_img = template.copy()
        draw = ImageDraw.Draw(test_img)
        for city, (lat, lon) in test_cities.items():
            x, y = self._latlon_to_map_xy(lat, lon, map_x, map_y, map_w, map_h)
            draw.ellipse([(x-10, y-10), (x+10, y+10)], fill=(255, 0, 0))
            draw.text((x+15, y-10), city, fill=(255, 0, 0))
        
        test_img.save('map_calibration_test.png')
        print("\nSaved test image to: map_calibration_test.png")
        print("Check if pins are in correct positions and adjust PAD_L, PAD_R, PAD_T, PAD_B values")
    
    def _parse_headshots(self, headshot_path):
        """
        Accepts:
          - single path string
          - comma-separated string "a.png,b.png"
          - list/tuple of paths
        Returns list[str]
        """
        if not headshot_path:
            return []
        if isinstance(headshot_path, (list, tuple)):
            return [p for p in headshot_path if p]
        if isinstance(headshot_path, str) and "," in headshot_path:
            return [p.strip() for p in headshot_path.split(",") if p.strip()]
        return [headshot_path]
    
    def _get_city_coordinates(self, city_name: str) -> tuple:
        """
        Get normalized coordinates (0-1) for US cities on a map.
        Returns (x, y) where x is 0 (west) to 1 (east), y is 0 (north) to 1 (south).
        FALLBACK ONLY - prefer using _geocode_city + _latlon_to_map_xy
        """
        # Normalize city name
        city_lower = city_name.lower().split(',')[0].strip()
        
        # Approximate US city positions (normalized coordinates) - Fallback only
        city_positions = {
            # West Coast
            'san francisco': (0.05, 0.42),
            'los angeles': (0.10, 0.60),
            'san diego': (0.10, 0.70),
            'seattle': (0.08, 0.20),
            'portland': (0.09, 0.25),
            # East Coast
            'new york': (0.92, 0.35),
            'boston': (0.95, 0.30),
            'philadelphia': (0.90, 0.38),
            'washington': (0.88, 0.40),
            'miami': (0.93, 0.80),
            'atlanta': (0.80, 0.60),
            # Central
            'chicago': (0.65, 0.35),
            'dallas': (0.48, 0.72),
            'houston': (0.50, 0.75),
            'austin': (0.48, 0.72),
            'denver': (0.40, 0.45),
            'phoenix': (0.25, 0.65),
            'las vegas': (0.20, 0.55),
            'detroit': (0.70, 0.32),
            'minneapolis': (0.55, 0.25),
            # Other major cities
            'san jose': (0.06, 0.44),
            'oakland': (0.05, 0.43),
            'sacramento': (0.07, 0.40),
        }
        
        # Try exact match first
        if city_lower in city_positions:
            return city_positions[city_lower]
        
        # Try partial matches
        for city, coords in city_positions.items():
            if city in city_lower or city_lower in city:
                return coords
        
        # Default to center of US
        return (0.50, 0.50)

    def create_slide(self, company_data: Dict, headshot_path: str, logo_path: str, map_path: Optional[str] = None, output_format: str = "pdf") -> bytes:
        # Check if template exists, try default locations if not found
        if not self.template_path or not os.path.exists(self.template_path):
            # Get the directory where this script is located
            script_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(script_dir) if os.path.basename(script_dir) in ['src', 'slauson-automation'] else script_dir
            
            # Try to find template in common locations (project root first)
            possible_paths = [
                # Project root (where user added the file)
                os.path.join(project_root, 'SLAUSON&CO.Template.pdf'),
                os.path.join(project_root, 'SLAUSON&CO.template'),
                os.path.join(project_root, 'SLAUSON&CO.Template'),
                # Templates folder
                os.path.join(project_root, 'templates', 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'templates', 'SLAUSON&CO.Template.pdf'),
                'templates/SLAUSON&CO.Template.pdf',
                # Script directory
                os.path.join(script_dir, 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'SLAUSON&CO.template'),
                # Current working directory
                'SLAUSON&CO.Template.pdf',
                'SLAUSON&CO.template',
                # Generic template names
                os.path.join(project_root, 'templates', 'template.pdf'),
                os.path.join(script_dir, 'templates', 'template.pdf'),
                'templates/template.pdf',
                'templates/template.png',
                'templates/template.jpg',
            ]
            found_path = None
            for path in possible_paths:
                if path and os.path.exists(path):
                    found_path = path
                    print(f"✓ Found template at: {found_path}")
                    break
            
            if not found_path:
                raise ValueError(
                    f"Template not found. Tried: {self.template_path}\n"
                    f"Please set SLIDE_TEMPLATE_PATH in .env or place template at project root: SLAUSON&CO.Template.pdf\n"
                    f"Script dir: {script_dir}, Project root: {project_root}\n"
                    f"Checked paths: {possible_paths[:5]}..."
                )
            self.template_path = found_path
        
        # Load template (PDF or image)
        template_ext = os.path.splitext(self.template_path)[1].lower()
        if template_ext == '.pdf':
            # Convert PDF to image
            template = self._pdf_to_image(self.template_path)
        else:
            # Load as image
            template = Image.open(self.template_path).convert('RGBA')
            template.load()  # Load fully before resize to avoid memory issues
        
        # Resize to standard slide size if needed
        if template.size != (1920, 1080):
            template = template.resize((1920, 1080), Image.Resampling.LANCZOS)
        
        # Create a copy to work with
        slide = template.copy()
        draw = ImageDraw.Draw(slide)
        width, height = slide.size
        
        # Load fonts (robust fallbacks for Render)
        name_font = self._load_font(140, bold=True)
        body_font = self._load_font(28)
        small_font = self._load_font(20)
        sidebar_font = self._load_font(28)
        
        # Extract company data
        company_name = company_data.get('name', '').upper()
        founders_text = company_data.get('founders', '')
        if isinstance(founders_text, list):
            founders_text = '\n'.join(founders_text)  # Use newlines instead of commas
        elif isinstance(founders_text, str) and ',' in founders_text:
            # Convert comma-separated to newline-separated
            founders_text = '\n'.join([f.strip() for f in founders_text.split(',')])
        
        co_investors_text = company_data.get('co_investors', '')
        if isinstance(co_investors_text, list):
            co_investors_text = '\n'.join(co_investors_text)  # Use newlines instead of commas
        elif isinstance(co_investors_text, str) and ',' in co_investors_text:
            # Convert comma-separated to newline-separated
            co_investors_text = '\n'.join([f.strip() for f in co_investors_text.split(',')])
        
        background_text = company_data.get('background', company_data.get('description', ''))
        location = company_data.get('address', company_data.get('location', 'Los Angeles'))
        if ',' in location:
            location = location.split(',')[0].strip()
        
        investment_stage = company_data.get('investment_stage', '')
        if not investment_stage:
            round_val = company_data.get('investment_round', 'PRE-SEED')
            quarter_val = company_data.get('quarter', 'Q2')
            year_val = company_data.get('year', '2024')
            investment_stage = f"{round_val} {quarter_val}, {year_val}"
        
        # 1. Replace company name (transparent background, significantly raised, very thick bold font, orange color)
        # Align with founders position
        founders_block_y = 400  # Approximate Y position of founders yellow block
        founders_text_x = 320  # Founders text X position
        
        # Detect orange US map bounding box from template (restricted to top-right ROI)
        # Do this once and reuse for both company name overlap detection and map pin placement
        map_area_x, map_area_y, map_width, map_height = self._detect_orange_us_bbox(template)
        
        # Company name position: aligned with founders but moved a bit to the left, significantly raised
        name_x = founders_text_x - 50  # Moved a bit to the left from founders position
        name_y = 120  # Lowered from 80 to move closer to content
        
        # Use orange color similar to the rest of the slide (lighter, more vibrant orange)
        name_color = (255, 140, 0)  # More vibrant orange, similar to slide orange
        
        # Dynamic font sizing based on company name length to prevent overlap with map
        # Start with base font size
        base_font_size = 180
        name_font_size = base_font_size
        
        # Calculate text width with base font to check for overlap
        test_font = self._load_font(name_font_size, bold=True)
        text_bbox = draw.textbbox((0, 0), company_name, font=test_font)
        text_width = text_bbox[2] - text_bbox[0]
        
        # If text would overlap with map (map starts at x=1250), reduce font size
        max_allowed_width = map_area_x - name_x - 50  # Leave 50px margin before map
        if text_width > max_allowed_width:
            # Calculate new font size to fit
            name_font_size = int((max_allowed_width / text_width) * base_font_size)
            # Ensure minimum readable size
            name_font_size = max(100, name_font_size)
            print(f"   Company name too long ({len(company_name)} chars), reducing font size to {name_font_size}px to avoid map overlap")
        
        # Use very thick, bold font (matching the image style - extra thick and bold)
        name_font = self._load_font(name_font_size, bold=True)
        
        # Draw company name with stroke (outline) to make it appear thicker
        # First draw the stroke (outline) in the same color but slightly darker
        stroke_color = (200, 80, 30)  # Slightly darker orange for stroke
        # Draw stroke by drawing text multiple times with slight offsets
        stroke_width = 2 if name_font_size > 150 else 1  # Adjust stroke based on font size
        for adj in range(-stroke_width, stroke_width + 1):
            for adj2 in range(-stroke_width, stroke_width + 1):
                if adj != 0 or adj2 != 0:
                    draw.text((name_x + adj, name_y + adj2), company_name, fill=stroke_color, font=name_font)
        
        # Then draw the main text on top
        draw.text((name_x, name_y), company_name, fill=name_color, font=name_font)
        
        # 2. Replace logo (circular, top right) - fit within circular bounds, higher position, transparent, more circular
        try:
            logo_img = Image.open(logo_path).convert('RGBA')
            # Load image fully before operations to avoid lazy loading issues
            logo_img.load()
            logo_x, logo_y = width - 170, 10  # Raised more (was 20, now 10) to avoid map overlap
            logo_size = 130
            
            # Don't erase background - keep it transparent (no black box)
            # Just paste the logo directly with circular mask
            
            # Resize logo to fit within circle (with more padding for circular shape)
            # Make it more circular by using square aspect ratio
            logo_img.thumbnail((logo_size - 20, logo_size - 20), Image.Resampling.LANCZOS)
            
            # Create perfectly circular mask
            circle_mask = Image.new('L', (logo_size, logo_size), 0)
            circle_draw = ImageDraw.Draw(circle_mask)
            circle_draw.ellipse([(0, 0), (logo_size, logo_size)], fill=255)
            
            # Center logo in circle (square crop for circular appearance)
            logo_w, logo_h = logo_img.size
            # Use the smaller dimension to make it more circular
            min_dim = min(logo_w, logo_h)
            logo_img = logo_img.crop(((logo_w - min_dim) // 2, (logo_h - min_dim) // 2, 
                                     (logo_w + min_dim) // 2, (logo_h + min_dim) // 2))
            # No need to resize again - thumbnail already resized it, just ensure exact size if needed
            if logo_img.size != (logo_size - 20, logo_size - 20):
            logo_img = logo_img.resize((logo_size - 20, logo_size - 20), Image.Resampling.LANCZOS)
            
            # Apply circular mask to logo
            logo_masked = Image.new('RGBA', (logo_size, logo_size), (0, 0, 0, 0))
            logo_masked.paste(logo_img, (10, 10), logo_img)  # Center with 10px padding
            logo_masked.putalpha(circle_mask)
            
            slide.paste(logo_masked, (logo_x, logo_y), logo_masked)
            draw = ImageDraw.Draw(slide)
        except Exception as e:
            print(f"Warning: Could not load logo: {e}")
        
        # 3. Updated Map Logic - Map is already in template, just add location text and adjust pin position
        try:
            # Reuse map bbox detected earlier (for company name overlap detection)
            # map_area_x, map_area_y, map_width, map_height already set above
            
            # Get lat/lon for city
            # Special case: Override Los Angeles to use adjusted coordinates (moved left)
            location_lower = location.lower().strip()
            if "los angeles" in location_lower or location_lower == "la":
                latlon = (34.0522, -118.5)  # Adjusted longitude to move pin left
            else:
                latlon = self._geocode_city(location) or self._fallback_latlon(location)
                if not latlon:
                    latlon = (39.5, -98.35)  # fallback center US
            
            # Debug prints to confirm bbox detection
            print(f"   DEBUG map bbox: ({map_area_x}, {map_area_y}, {map_width}, {map_height})")
            print(f"   DEBUG location: {location}, latlon: {latlon}")
            
            # Note: AK (Juneau) and HI (Honolulu) coordinates will be clamped to contiguous US bounds
            # in _latlon_to_map_xy, so they won't appear accurately on the map. Consider skipping
            # pin placement for these states or implementing inset logic if needed.
            lat, lon = latlon
            pin_x, pin_y = self._latlon_to_map_xy(lat, lon, map_area_x, map_area_y, map_width, map_height)
            # REMOVED: pin_y -= 6  # Remove aesthetic offset as it causes inaccuracy
            
            print(f"   PIN: ({pin_x}, {pin_y})")
            
            # Draw the pin (yellow location pin icon - teardrop shape with circular hole)
            yellow = (255, 215, 0)
            black = (0, 0, 0)
            
            # Pin dimensions - larger for visibility
            pin_width = 40  # Width of the rounded top
            pin_height = 50  # Total height including point
            point_length = 12  # Length of the pointed bottom
            hole_radius = 6  # Radius of the circular hole
            
            # Create a temporary image for the pin to draw the teardrop shape
            pin_img_size = max(pin_width, pin_height + point_length) + 20
            pin_img = Image.new('RGBA', (pin_img_size, pin_img_size), (0, 0, 0, 0))
            pin_draw = ImageDraw.Draw(pin_img)
            
            # Calculate center position in the pin image
            center_x = pin_img_size // 2
            center_y = pin_img_size // 2
            
            # Draw shadow first (slightly offset and darker) for depth
            shadow_offset = 3
            shadow_center_x = center_x + shadow_offset
            shadow_center_y = center_y + shadow_offset
            
            # Shadow: rounded top (ellipse) + pointed bottom (triangle)
            # Top ellipse for shadow
            shadow_top_y = shadow_center_y - pin_height // 2
            pin_draw.ellipse(
                [(shadow_center_x - pin_width // 2, shadow_top_y - pin_width // 4),
                 (shadow_center_x + pin_width // 2, shadow_top_y + pin_width // 4)],
                fill=(50, 50, 50, 120)
            )
            # Bottom triangle for shadow
            shadow_points = [
                (shadow_center_x - pin_width // 2, shadow_top_y + pin_width // 4),
                (shadow_center_x + pin_width // 2, shadow_top_y + pin_width // 4),
                (shadow_center_x, shadow_center_y + point_length + shadow_offset),
            ]
            pin_draw.polygon(shadow_points, fill=(50, 50, 50, 120))
            
            # Draw main teardrop shape (yellow)
            # Top: rounded ellipse (the rounded part of the teardrop)
            top_y = center_y - pin_height // 2
            pin_draw.ellipse(
                [(center_x - pin_width // 2, top_y - pin_width // 4),
                 (center_x + pin_width // 2, top_y + pin_width // 4)],
                fill=yellow, outline=black, width=2
            )
            
            # Bottom: triangle connecting to point (the pointed part)
            teardrop_points = [
                (center_x - pin_width // 2, top_y + pin_width // 4),  # Left bottom of ellipse
                (center_x + pin_width // 2, top_y + pin_width // 4),  # Right bottom of ellipse
                (center_x, center_y + point_length),  # Point at bottom
            ]
            pin_draw.polygon(teardrop_points, fill=yellow, outline=black, width=2)
            
            # Draw circular hole in the center (dark circle for depth)
            hole_center_x = center_x
            hole_center_y = top_y  # Position hole at top center
            pin_draw.ellipse(
                [(hole_center_x - hole_radius, hole_center_y - hole_radius),
                 (hole_center_x + hole_radius, hole_center_y + hole_radius)],
                fill=black
            )
            
            # Add inner highlight ring for 3D effect (lighter inner ring)
            pin_draw.ellipse(
                [(hole_center_x - hole_radius + 2, hole_center_y - hole_radius + 2),
                 (hole_center_x + hole_radius - 2, hole_center_y + hole_radius - 2)],
                fill=(100, 100, 100)
            )
            
            # Paste the pin onto the main slide so the bottom tip lands on (pin_x, pin_y)
            paste_pin_x = pin_x - pin_img_size // 2
            paste_pin_y = pin_y - pin_img_size // 2
            # Move image up so the bottom tip hits (pin_x, pin_y) instead of centering
            paste_pin_y -= int(pin_img_size * 0.20)  # Adjust 0.18-0.25 if needed
            slide.paste(pin_img, (paste_pin_x, paste_pin_y), pin_img)
            draw = ImageDraw.Draw(slide)
            
            # Place the city name label with yellow block (similar to other yellow blocks)
            yellow = (255, 215, 0)  # Yellow for the block
            black = (0, 0, 0)  # Black text on yellow background
            
            # Use larger, extra-bold appearance for location text
            location_font = self._load_font(32, bold=True)
            location_stroke_width = 1
            
            # Get text dimensions with larger font
            label_bbox = draw.textbbox((0, 0), location, font=location_font, stroke_width=location_stroke_width)
            text_width = label_bbox[2] - label_bbox[0]
            text_height = label_bbox[3] - label_bbox[1]
            
            # Yellow box dimensions (bigger padding for larger text)
            box_padding_x = 24  # Increased from 20
            box_padding_y = 12  # Increased from 10
            box_width = text_width + box_padding_x * 2
            box_height = text_height + box_padding_y * 2
            
            # Position box with more space from pin; raise significantly to avoid being too low
            box_x = pin_x + 40  # More space from pin (was 15, now 40)
            box_y = pin_y - 70  # Raised more (was -50, now -70 to move label higher)
            
            # Make sure box stays within map bounds
            if box_x + box_width > map_area_x + map_width:
                box_x = pin_x - box_width - 40  # Move to the left if it would go outside
            if box_y + box_height > map_area_y + map_height:
                box_y = pin_y - box_height - 20  # Move up if it would go outside
            
            # Draw yellow box (similar to other section labels)
            draw.rectangle([(box_x, box_y), (box_x + box_width, box_y + box_height)], fill=yellow)
            
            # Draw location text in the yellow box (centered with padding)
            text_x = box_x + box_padding_x
            text_y = box_y + box_padding_y
            draw.text(
                (text_x, text_y),
                location,
                fill=black,
                font=location_font,
                stroke_width=location_stroke_width,
                stroke_fill=black
            )
            
        except Exception as e:
            print(f"Warning: Map update failed: {e}")
        
        # 4. Replace Founders text (moved left and down, transparent)
        # Detect yellow block area for founders (left side, below company name)
        founders_block_y = 400  # Approximate Y position of founders yellow block
        founders_text_y = founders_block_y + 15  # Moved down a little
        founders_text_x = 320  # Moved a little to the left (was 350)
        founders_color = self._get_text_color_from_template(template, founders_text_x, founders_text_y, 600, 100)
        
        # Don't erase background - keep it transparent (no black box)
        # Draw founders text with newlines directly
        founders_lines = founders_text.split('\n')
        for i, line in enumerate(founders_lines):
            draw.text((founders_text_x, founders_text_y + i * 35), line, fill=founders_color, font=body_font)
        
        # 5. Replace Co-Investors text (separate text box, to the right of founders, aligned with founders height)
        investors_block_y = 500  # Approximate Y position of co-investors yellow block
        investors_text_y = founders_text_y  # Aligned with founders (same height)
        investors_text_x = 650  # Moved to the right of founders (founders is at 320, so co-investors at 650)
        investors_color = self._get_text_color_from_template(template, investors_text_x, investors_text_y, 600, 100)
        
        # Don't erase background - keep it transparent (no black box)
        # Draw co-investors text with newlines directly (separate text box)
        investors_lines = co_investors_text.split('\n')
        for i, line in enumerate(investors_lines):
            draw.text((investors_text_x, investors_text_y + i * 35), line, fill=investors_color, font=body_font)
        
        # 6. Replace Background text (aligned with founders, wider text area, bigger font, transparent background)
        bg_block_y = 600  # Approximate Y position of background yellow block
        bg_text_y = bg_block_y + 50  # Text starts below yellow block
        bg_text_x = founders_text_x  # Aligned with founders (moved left from 400 to match founders at 320)
        bg_text_width = 700  # Wider text area (was 500) to fill space more
        words = background_text.split()
        lines = []
        current_line = []
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=body_font)
            if bbox[2] - bbox[0] < bg_text_width:  # Wider width for background
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        if current_line:
            lines.append(' '.join(current_line))
        
        # Get text color from template (don't erase background - keep it transparent)
        bg_color = self._get_text_color_from_template(template, bg_text_x, bg_text_y, bg_text_width, 200)
        
        # Use slightly bigger font for background text
        bg_font = self._load_font(32)
        
        # Draw text directly without erasing background (transparent)
        for i, line in enumerate(lines[:10]):
            draw.text((bg_text_x, bg_text_y + i * 36), line, fill=bg_color, font=bg_font)  # Slightly more spacing
        
        # 7. Replace headshots (below map, moved left, bigger size, transparent)
        headshot_processed = False
        try:
            from image_processor import ImageProcessor
            
            # First, validate that the headshot file exists and is a valid image
            # If not, skip headshot processing (don't crash the entire slide generation)
            if not headshot_path or not os.path.exists(headshot_path):
                print(f"Warning: Headshot file not found: {headshot_path}, skipping headshot")
                headshot_path = None
            
            if headshot_path:
                # Try to open the original image first to validate it (quick check)
                try:
                    test_img = Image.open(headshot_path)
                    # Don't verify - just check if we can open it (verify is slow)
                    test_img.close()
                except Exception as e:
                    print(f"Warning: Headshot file is not a valid image: {e}, skipping headshot")
                    headshot_path = None
            
            headshot_paths = self._parse_headshots(headshot_path)
            headshot_paths = [p for p in headshot_paths if p and os.path.exists(p)]
            if not headshot_paths:
                print("   Skipping headshot processing (no valid headshot file)")
            else:
                # Remove background from headshot(s) to make them transparent
                use_api_removal = False
                try:
                    from config import Config
                    if hasattr(Config, 'REMOVEBG_API_KEY') and Config.REMOVEBG_API_KEY and Config.REMOVEBG_API_KEY.strip():
                        use_api_removal = True
                except:
                    pass
                
                print("   REMOVEBG_API_KEY not set, using manual background removal" if not use_api_removal else "   Using remove.bg API when available")

                # Headshot target box (tuned to match reference slide)
                headshot_area_width = 820
                headshot_area_height = 900
                headshot_area_x = map_area_x + (map_width - headshot_area_width) // 2 + 20
                headshot_area_y = map_area_y + int(map_height * 0.46)

                def load_process_headshot(path: str) -> Optional[Image.Image]:
                    try:
                        img = self._remove_bg_best_effort(path, use_api_removal)

                        # Tiny edge soften after we have a good alpha
                        a = img.split()[-1].filter(ImageFilter.GaussianBlur(radius=0.6))
                        img.putalpha(a)

                        # Convert to greyscale while preserving alpha
                        r, g, b, a = img.split()
                        gray = img.convert('L')
                        img = Image.merge('RGBA', (gray, gray, gray, a))
                        return img
                except Exception as e:
                        print(f"Warning: failed headshot {path}: {e}")
                        return None

                imgs = [load_process_headshot(p) for p in headshot_paths[:2]]
                imgs = [im for im in imgs if im is not None]

                if len(imgs) == 1:
                    im = imgs[0]
                    im.thumbnail((headshot_area_width, headshot_area_height), Image.Resampling.LANCZOS)
                    paste_x = headshot_area_x + (headshot_area_width - im.size[0]) // 2
                    paste_y = headshot_area_y + (headshot_area_height - im.size[1]) // 2
                    slide.paste(im, (paste_x, paste_y), im.split()[3])
                    draw = ImageDraw.Draw(slide)
                elif len(imgs) == 2:
                    gap = 30
                    each_w = (headshot_area_width - gap) // 2
                    each_h = headshot_area_height
                    left, right = imgs
                    left.thumbnail((each_w, each_h), Image.Resampling.LANCZOS)
                    right.thumbnail((each_w, each_h), Image.Resampling.LANCZOS)

                    base_y = headshot_area_y + headshot_area_height
                    left_x = headshot_area_x + (each_w - left.size[0]) // 2
                    right_x = headshot_area_x + each_w + gap + (each_w - right.size[0]) // 2
                    left_y = base_y - left.size[1]
                    right_y = base_y - right.size[1]

                    slide.paste(left, (left_x, left_y), left.split()[3])
                    slide.paste(right, (right_x, right_y), right.split()[3])
                draw = ImageDraw.Draw(slide)
        except Exception as e:
            print(f"Warning: Could not load headshot: {e}")
            import traceback
            traceback.print_exc()
        
        # 8. Replace investment stage in sidebar (transparent, black, thicker text like company name, includes year)
        # Find Slauson&Co position in sidebar (bottom of sidebar)
        slauson_y = height - 200  # Approximate position of "SLAUSON&CO." text (bottom of sidebar)
        
        # Parse investment stage to format as "STAGE QUARTER YEAR" (e.g., "SEED Q2 2024")
        # Format is typically like "SEED Q2, 2024" or "PRE-SEED Q2, 2024"
        stage_parts = investment_stage.upper().split(',')
        if len(stage_parts) >= 2:
            stage_quarter = stage_parts[0].strip()  # "SEED Q2" or "PRE-SEED Q2"
            year_text = stage_parts[1].strip()  # "2024"
            # Combine into "STAGE QUARTER YEAR" format
            stage_text = f"{stage_quarter} {year_text}"  # "SEED Q2 2024"
        else:
            stage_text = investment_stage.upper()
        
        # Use bigger font size for better visibility (36-46pt range)
        stage_font_size = 40
        sidebar_bold_font = self._load_font(stage_font_size, bold=True)
        
        # Don't erase background - keep it transparent
        # Use black color
        stage_color = (0, 0, 0)  # Black color
        stroke_color = (50, 50, 50)  # Dark grey for stroke (thicker effect)
        
        # Draw investment stage at the TOP of sidebar (opposite end from Slauson&Co)
        # Align with Slauson&Co position (which is at the bottom of sidebar, left-aligned)
        # Slauson&Co is positioned more to the left, around x=10-20 in the sidebar
        sidebar_width = 200
        slauson_x = 10  # More to the left to align with Slauson&Co text
        
        stage_top_y = 80  # Near the top of the sidebar
        
        # Get text dimensions first to calculate proper image size
        temp_img = Image.new('RGBA', (1000, 1000), (0, 0, 0, 0))
        temp_draw = ImageDraw.Draw(temp_img)
        bbox = temp_draw.textbbox((0, 0), stage_text, font=sidebar_bold_font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        # Create image for single-line text, rotated -90 degrees (same as SLAUSON&CO)
        # When rotated -90 degrees, width becomes height and height becomes width
        padding = 60  # Reduced padding to prevent excessive scaling
        stage_img_width = text_height + padding * 2  # Width after rotation
        stage_img_height = text_width + padding * 2  # Height after rotation
        
        # Ensure minimum size
        stage_img_width = max(stage_img_width, 400)
        stage_img_height = max(stage_img_height, 700)
        
        stage_img = Image.new('RGBA', (stage_img_width, stage_img_height), (0, 0, 0, 0))
        stage_draw = ImageDraw.Draw(stage_img)
        
        # Draw text normally (no word reversal) at the top of the image
        # When rotated -90 degrees, the top becomes the right side
        # Position text at top so after rotation it's properly positioned
        # Center text horizontally before rotation
        text_x = stage_img_width // 2 - text_width // 2
        text_y = padding  # Position at top edge - becomes right edge after rotation
        
        # Draw stroke by drawing text multiple times with slight offsets (thicker effect)
        for adj in range(-3, 4):  # slightly less bold
            for adj2 in range(-3, 4):
                if adj != 0 or adj2 != 0:
                    stage_draw.text((text_x + adj, text_y + adj2), stage_text, fill=stroke_color, font=sidebar_bold_font)
        
        # Then draw the main text on top
        stage_draw.text((text_x, text_y), stage_text, fill=stage_color, font=sidebar_bold_font)
        
        # Rotate so text reads bottom -> top (same as SLAUSON&CO)
        stage_img = stage_img.rotate(90, expand=True)
        
        # --- Fit the rotated image into the sidebar safely ---
        sidebar_w = 200
        top_margin = 55  # Adjusted for better fit
        bottom_margin = 15  # Reduced to give more usable height
        side_margin = 6  # Reduced to give more usable width
        
        max_w = sidebar_w - 2 * side_margin
        max_h = height - top_margin - bottom_margin  # height is slide height (1080)
        
        final_w, final_h = stage_img.size
        
        # Scale DOWN if needed to fit, but don't let it get microscopic
        scale = min(max_w / final_w, max_h / final_h, 1.0)
        scale = max(scale, 0.85)  # Minimum scale floor to prevent tiny text
        if scale < 1.0:
            new_w = max(1, int(final_w * scale))
            new_h = max(1, int(final_h * scale))
            stage_img = stage_img.resize((new_w, new_h), Image.Resampling.LANCZOS)
            final_w, final_h = stage_img.size
        
        # Compute placement
        paste_x = 12              # align with SLAUSON&CO left
        paste_y = top_margin + 25 # lower it
        
        # Clamp so it can NEVER go off-canvas
        paste_x = max(0, min(paste_x, sidebar_w - final_w))
        paste_y = max(0, min(paste_y, height - final_h))
        
        slide.paste(stage_img, (paste_x, paste_y), stage_img)
        
        # Convert to RGB for PDF
        slide_rgb = Image.new('RGB', slide.size, (42, 42, 42))
        slide_rgb.paste(slide, mask=slide.split()[3] if slide.mode == 'RGBA' else None)
        
        # Return based on output format
        if output_format.lower() == "pptx":
            return self._create_pptx_from_slide(
                slide, company_data, headshot_path, logo_path, map_path,
                map_area_x, map_area_y, map_width, map_height
            )
        
        # Default: Convert to PDF (flattened, not editable)
        # Convert to RGB for PDF
        slide_rgb = Image.new('RGB', slide.size, (42, 42, 42))
        slide_rgb.paste(slide, mask=slide.split()[3] if slide.mode == 'RGBA' else None)
        
        # Convert to PDF
        # Use temporary file to avoid PIL's fileno() issue with BytesIO
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            try:
                slide_rgb.save(tmp_file.name, format='PNG', optimize=False)
                with open(tmp_file.name, 'rb') as f:
                    img_bytes = f.read()
                pdf_bytes = img2pdf.convert(img_bytes)
        return pdf_bytes
            finally:
                # Clean up temporary file
                try:
                    os.unlink(tmp_file.name)
                except:
                    pass
    
    def _create_pptx_from_slide(
        self, 
        template_slide: Image.Image,
        company_data: Dict,
        headshot_path: str,
        logo_path: str,
        map_path: Optional[str],
        map_area_x: int,
        map_area_y: int,
        map_width: int,
        map_height: int
    ) -> bytes:
        """
        Create an editable PPTX file with separate elements for Canva import.
        This allows headshot, logo, map, and text to be editable in Canva.
        
        Args:
            template_slide: PIL Image of the template background
            company_data: Company information
            headshot_path: Path to headshot image
            logo_path: Path to logo image
            map_path: Path to map image (optional)
            map_area_x, map_area_y, map_width, map_height: Map position/size
            
        Returns:
            PPTX file bytes
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX generation. Install with: pip install python-pptx")
        
        # Slide dimensions: 1920x1080 pixels = 20x11.25 inches at 96 DPI
        # Standard widescreen: 13.333" x 7.5" (16:9)
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 1920px / 144 DPI
        prs.slide_height = Inches(7.5)   # 1080px / 144 DPI
        
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide_pptx = prs.slides.add_slide(slide_layout)
        
        # CRITICAL: Don't add full-slide background image - it causes Canva to flatten everything into one image
        # Instead, we'll add ONLY the editable elements (logo, headshot, text, map) on a blank slide
        # The user can add their own background in Canva, or we can add it as a separate editable element later
        # For now, just add a simple colored background using shapes (editable)
        from pptx.enum.shapes import MSO_SHAPE
        
        # 1) Orange sidebar (left side, editable shape)
        sidebar_width_inch = Inches(200 / 96.0)  # 200px converted to inches
        sidebar = slide_pptx.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            sidebar_width_inch,
            prs.slide_height
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = RGBColor(242, 140, 40)  # Orange #F28C28
        sidebar.line.fill.background()  # No border
        
        # 2) Dark grey main area (right side, editable shape)
        main_area = slide_pptx.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            sidebar_width_inch,
            Inches(0),
            prs.slide_width - sidebar_width_inch,
            prs.slide_height
        )
        main_area.fill.solid()
        main_area.fill.fore_color.rgb = RGBColor(42, 42, 42)  # Dark grey
        main_area.line.fill.background()  # No border
        
        # Get actual positions from template (convert pixels to inches at 96 DPI)
        width, height = template_slide.size
        px_to_inch = 1.0 / 96.0  # 96 DPI standard
        
        # Add "SLAUSON&CO." text in sidebar (bottom, editable text)
        slauson_text = "SLAUSON&CO."
        slauson_y_px = height - 200  # Bottom of sidebar
        slauson_x_px = 10  # Left side of sidebar
        
        slauson_textbox = slide_pptx.shapes.add_textbox(
            Inches(slauson_x_px * px_to_inch),
            Inches(slauson_y_px * px_to_inch),
            Inches(180 * px_to_inch),  # Width for sidebar
            Inches(1.5)
        )
        slauson_tf = slauson_textbox.text_frame
        slauson_tf.word_wrap = False
        slauson_p = slauson_tf.paragraphs[0]
        slauson_run = slauson_p.add_run()
        slauson_run.text = slauson_text
        slauson_run.font.size = Pt(28)
        slauson_run.font.bold = True
        slauson_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # 2) Logo (top right, editable) - from PIL code: logo_x, logo_y = width - 170, 10
        if logo_path and os.path.exists(logo_path):
            logo_img = Image.open(logo_path).convert('RGBA')
            # Load image fully before operations to avoid lazy loading issues
            logo_img.load()
            # Make circular (same as PIL code)
            logo_size_px = 130
            logo_img.thumbnail((logo_size_px - 20, logo_size_px - 20), Image.Resampling.LANCZOS)
            
            # Create circular mask
            circle_mask = Image.new('L', (logo_size_px, logo_size_px), 0)
            circle_draw = ImageDraw.Draw(circle_mask)
            circle_draw.ellipse([(0, 0), (logo_size_px, logo_size_px)], fill=255)
            
            # Apply mask
            logo_w, logo_h = logo_img.size
            min_dim = min(logo_w, logo_h)
            logo_img = logo_img.crop(((logo_w - min_dim) // 2, (logo_h - min_dim) // 2, 
                                     (logo_w + min_dim) // 2, (logo_h + min_dim) // 2))
            # Only resize if thumbnail didn't produce exact size
            if logo_img.size != (logo_size_px - 20, logo_size_px - 20):
            logo_img = logo_img.resize((logo_size_px - 20, logo_size_px - 20), Image.Resampling.LANCZOS)
            logo_masked = Image.new('RGBA', (logo_size_px, logo_size_px), (0, 0, 0, 0))
            logo_masked.paste(logo_img, (10, 10), logo_img)
            logo_masked.putalpha(circle_mask)
            
            logo_bytes = io.BytesIO()
            logo_masked.save(logo_bytes, format='PNG')
            logo_bytes.seek(0)
            
            # Position: from PIL code - logo_x, logo_y = width - 170, 10
            logo_x_inch = (width - 170) * px_to_inch
            logo_y_inch = 10 * px_to_inch
            logo_size_inch = logo_size_px * px_to_inch
            
            slide_pptx.shapes.add_picture(
                logo_bytes,
                Inches(logo_x_inch),
                Inches(logo_y_inch),
                width=Inches(logo_size_inch)
            )
        
        # 3) Company name (editable text) - from PIL code: name_x, name_y = founders_text_x - 50, 120
        company_name = company_data.get("name", "").upper()
        if company_name:
            # Position from PIL: name_x = founders_text_x - 50, name_y = 120
            # founders_text_x = 320 from PIL code
            name_x_px = 320 - 50  # 270
            name_y_px = 120
            
            # Calculate font size (same logic as PIL)
            base_font_size = 180
            name_font_size = base_font_size
            # Check for overlap with map
            max_allowed_width = map_area_x - name_x_px - 50
            # Estimate text width (rough calculation)
            estimated_width = len(company_name) * (name_font_size * 0.6)  # Rough estimate
            if estimated_width > max_allowed_width:
                name_font_size = int((max_allowed_width / estimated_width) * base_font_size)
                name_font_size = max(100, name_font_size)
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(name_x_px * px_to_inch),
                Inches(name_y_px * px_to_inch),
                Inches((map_area_x - name_x_px - 50) * px_to_inch),
                Inches(2.0)  # Height for large text
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = company_name
            run.font.size = Pt(name_font_size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 140, 0)  # Orange color
        
        # 4) Investment stage (editable text) - positioned in sidebar
        investment_stage = company_data.get("investment_stage", "")
        if not investment_stage:
            round_val = company_data.get("investment_round", "PRE-SEED")
            quarter_val = company_data.get("quarter", "Q2")
            year_val = company_data.get("year", "2024")
            investment_stage = f"{round_val} {quarter_val} {year_val}"
        
        if investment_stage:
            paste_x_px = 15   # slightly right
            paste_y_px = 70   # lower
            
            # Make the textbox TALL/WIDE BEFORE rotation so it has room after rotation.
            # Use most of the slide height.
            box_w_px = 180
            box_h_px = 950   # <-- big, prevents "tiny" rendering
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(paste_x_px * px_to_inch),
                Inches(paste_y_px * px_to_inch),
                Inches(box_w_px * px_to_inch),
                Inches(box_h_px * px_to_inch),
            )
            textbox.rotation = 270  # bottom->top like the sidebar
            
            tf = textbox.text_frame
            tf.word_wrap = False
            
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = investment_stage
            
            # BIGGER FONT (36-46pt range)
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # 5) Headshot (below map, editable)
        if headshot_path and os.path.exists(headshot_path):
            headshot_bytes = io.BytesIO()
            headshot_img = Image.open(headshot_path).convert('RGBA')
            headshot_img.load()  # Load fully before save to avoid memory issues
            # Ensure transparent background for PPTX as well (grayscale-aware)
            try:
                headshot_img = self._remove_background_gray(headshot_img, tol=12, feather=1)
            except Exception as e:
                print(f"   Warning: PPTX fallback background removal failed: {e}")
            # Process headshot (background removal already done in PIL code)
            headshot_img.save(headshot_bytes, format='PNG', optimize=False)
            headshot_bytes.seek(0)
            
            # Position from PIL: headshot_area_width = 550 * 2.2 = 1210, headshot_area_height = 500 * 2.2 = 1100
            # headshot_area_x = map_area_x + (map_width - headshot_area_width) // 2 - 50
            # headshot_area_y = map_area_y + map_height - 50
            # Headshot target box (tuned to match reference slide)
            headshot_area_width_px = 820
            headshot_area_height_px = 900
            headshot_area_x_px = map_area_x + (map_width - headshot_area_width_px) // 2 + 20
            headshot_area_y_px = map_area_y + int(map_height * 0.46)
            
            slide_pptx.shapes.add_picture(
                headshot_bytes,
                Inches(headshot_area_x_px * px_to_inch),
                Inches(headshot_area_y_px * px_to_inch),
                width=Inches(headshot_area_width_px * px_to_inch),
                height=Inches(headshot_area_height_px * px_to_inch)
            )
        
        # 6) Founders text (editable) - from PIL code: founders_text_x = 320, founders_text_y = 415
        founders_text = company_data.get("founders", "")
        if founders_text:
            if isinstance(founders_text, list):
                founders_text = '\n'.join(founders_text)
            elif isinstance(founders_text, str) and ',' in founders_text:
                founders_text = '\n'.join([f.strip() for f in founders_text.split(',')])
            
            # Position from PIL: founders_text_x = 320, founders_text_y = 415
            founders_text_x_px = 320
            founders_text_y_px = 415
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(founders_text_x_px * px_to_inch),
                Inches(founders_text_y_px * px_to_inch),
                Inches(3.0),
                Inches(2.0)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = founders_text
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # 7) Co-investors text (editable) - from PIL code: investors_text_x = 650, investors_text_y = 415
        co_investors_text = company_data.get("co_investors", "")
        if co_investors_text:
            if isinstance(co_investors_text, list):
                co_investors_text = '\n'.join(co_investors_text)
            elif isinstance(co_investors_text, str) and ',' in co_investors_text:
                co_investors_text = '\n'.join([f.strip() for f in co_investors_text.split(',')])
            
            # Position from PIL: investors_text_x = 650, investors_text_y = 415
            investors_text_x_px = 650
            investors_text_y_px = 415
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(investors_text_x_px * px_to_inch),
                Inches(investors_text_y_px * px_to_inch),
                Inches(3.0),
                Inches(2.0)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = co_investors_text
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # 8) Background text (editable) - from PIL code: bg_text_x = 320, bg_text_y = 650
        background_text = company_data.get("background", company_data.get("description", ""))
        if background_text:
            # Position from PIL: bg_text_x = 320, bg_text_y = 650, bg_text_width = 700
            bg_text_x_px = 320
            bg_text_y_px = 650
            bg_text_width_px = 700
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(bg_text_x_px * px_to_inch),
                Inches(bg_text_y_px * px_to_inch),
                Inches(bg_text_width_px * px_to_inch),
                Inches(2.5)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = background_text
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # 9) Map (if provided, editable) - positioned at detected map area
        if map_path and os.path.exists(map_path):
            map_bytes = io.BytesIO()
            map_img = Image.open(map_path).convert('RGBA')
            map_img.load()  # Load fully before save to avoid memory issues
            map_img.save(map_bytes, format='PNG', optimize=False)
            map_bytes.seek(0)
            
            # Position: top right area (from detected map area)
            map_x = map_area_x / 96.0  # Convert pixels to inches
            map_y = map_area_y / 96.0
            map_w = map_width / 96.0
            map_h = map_height / 96.0
            
            slide_pptx.shapes.add_picture(
                map_bytes,
                Inches(map_x),
                Inches(map_y),
                width=Inches(map_w),
                height=Inches(map_h)
            )
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        return pptx_bytes.getvalue()
